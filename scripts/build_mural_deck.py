"""Generate an editable slide deck of the mural mockup from mural-mockup.json.

One source of truth: this reads data/mural-mockup.json (built by
build_mural_mockup.py) and produces a .pptx that embeds the actual preview
images as bytes — no image hosting required. Uploaded to Google Drive it opens
as an editable Google Slides deck, so Kelly can co-edit and present from Drive.
Rerun after re-sourcing to regenerate a fresh deck.

Deck: title → how-to-read → one slide per frame (chronological, gaps grayscaled
and flagged "To source" with their track) → closing "what we need" ask.

Output: data/research-assets/beaumont-mural-mockup.pptx
"""
import io
import json
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

ROOT = Path(__file__).resolve().parent.parent
DATA = ROOT / "data" / "mural-mockup.json"
SITE = ROOT / "site"
OUT = ROOT / "data" / "research-assets" / "beaumont-mural-mockup.pptx"

# Site palette.
INK = RGBColor(0x1E, 0x28, 0x23)
PAPER = RGBColor(0xF4, 0xF0, 0xE8)
PANEL = RGBColor(0xFF, 0xFD, 0xF8)
FOREST = RGBColor(0x24, 0x3C, 0x31)
FOREST2 = RGBColor(0x36, 0x54, 0x47)
RUST = RGBColor(0xA6, 0x49, 0x2D)
GOLD = RGBColor(0xD3, 0xA5, 0x4D)
MUTED = RGBColor(0x68, 0x73, 0x6D)
READY = RGBColor(0x2F, 0x6B, 0x4F)
TRACK_COLOR = {"A": RGBColor(0x2D, 0x6E, 0xA6), "B": RGBColor(0x8A, 0x5A, 0x1F),
               "C": RGBColor(0x7A, 0x4B, 0xB0)}
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SERIF = "Georgia"
SANS = "Calibri"

EMU_IN = 914400
W, H = Inches(13.333), Inches(7.5)


def solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, x, y, w, h, color):
    from pptx.enum.shapes import MSO_SHAPE
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    solid(s, color)
    s.shadow.inherit = False
    return s


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         wrap=True):
    """runs: list of paragraphs; each paragraph is list of (text, font, size, color, bold, italic)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = Pt(0)
    tf.margin_top = tf.margin_bottom = Pt(0)
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        for (t, font, size, color, bold, italic) in para:
            r = p.add_run()
            r.text = t
            r.font.name = font
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.italic = italic
    return tb


def chip(slide, x, y, label, fill, txt=WHITE):
    w = Inches(0.14 + 0.092 * len(label))
    r = rect(slide, x, y, w, Inches(0.34), fill)
    tf = r.text_frame
    tf.margin_left = tf.margin_right = Pt(6)
    tf.margin_top = tf.margin_bottom = Pt(2)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = label
    run.font.name = SANS; run.font.size = Pt(11); run.font.bold = True
    run.font.color.rgb = txt
    return w


def blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def load_img(frame):
    """Return (BytesIO, w_px, h_px). Grayscale gap frames to match the web mockup."""
    rel = frame.get("preview") or frame.get("thumbnail")
    p = SITE / rel
    im = Image.open(p).convert("RGB")
    if frame["kind"] == "gap":
        im = im.convert("L").point(lambda v: int(v * 0.92 + 8)).convert("RGB")
    # Cap size — slides never need more than ~1400px on the long edge.
    long_edge = max(im.width, im.height)
    if long_edge > 1400:
        scale = 1400 / long_edge
        im = im.resize((round(im.width * scale), round(im.height * scale)), Image.LANCZOS)
    buf = io.BytesIO()
    im.save(buf, format="JPEG", quality=88)
    buf.seek(0)
    return buf, im.width, im.height


def fit(w_px, h_px, box_w, box_h):
    """Scale to fit inside box (contain), return (w_emu, h_emu)."""
    scale = min(box_w / w_px, box_h / h_px)
    return Emu(int(w_px * scale)), Emu(int(h_px * scale))


def title_slide(prs, data):
    s = blank(prs)
    rect(s, 0, 0, W, H, FOREST)
    rect(s, 0, 0, W, Inches(0.14), GOLD)
    text(s, Inches(0.9), Inches(1.1), Inches(11.5), Inches(0.5),
         [[("TIMELINE MURAL · WORKING PROOF", SANS, 14, GOLD, True, False)]])
    text(s, Inches(0.85), Inches(1.6), Inches(11.6), Inches(1.6),
         [[("Beaumont, in one long wall", SERIF, 54, WHITE, False, False)]])
    text(s, Inches(0.9), Inches(3.2), Inches(9.8), Inches(1.4),
         [[("The mural we can print today — production-grade stand-ins for "
            "Kelly's vision — laid end to end with the images still to be sourced "
            "shown in place, so the whole story arc is visible and the gaps are "
            "honest.", SERIF, 19, RGBColor(0xDC, 0xE4, 0xDE), False, True)]])
    stats = [(str(data["ready_count"]), "ready to print now"),
             (str(data["gap_count"]), "to re-acquire from source"),
             (_span(data), "years spanned")]
    x = Inches(0.9)
    for value, label in stats:
        text(s, x, Inches(5.0), Inches(3.6), Inches(1.2),
             [[(value, SERIF, 40, GOLD, False, False)],
              [(label, SANS, 13, RGBColor(0xB9, 0xC6, 0xBF), False, False)]])
        x = Emu(x + Inches(3.7))
    return s


def howto_slide(prs):
    s = blank(prs)
    rect(s, 0, 0, W, H, PAPER)
    rect(s, 0, 0, Inches(0.22), H, GOLD)
    text(s, Inches(0.9), Inches(0.8), Inches(11.5), Inches(0.9),
         [[("How to read this deck", SERIF, 34, INK, False, False)]])
    lines = [
        ("Each slide is one frame of the mural, in chronological order — walk them "
         "front to back and you walk Beaumont from before the town to today.", INK),
        ("READY frames are production-grade images we can print now, chosen to match "
         "Kelly's curated set. They appear in full colour with their print size.", READY),
        ("TO SOURCE frames are Kelly's picks that survive only as low-resolution "
         "scans. They appear in grayscale, flagged with where to re-acquire a better "
         "original: Track A online, Track B Malki Museum + Morongo Tribe, Track C "
         "Library storage.", RUST),
        ("The caption under each image is written for a visitor to the finished wall. "
         "The case for why the image belongs in the timeline is in each slide's "
         "speaker notes — your talking points, not the audience's. Captions are draft "
         "until facts and rights are confirmed.", MUTED),
    ]
    y = Inches(2.1)
    for body, color in lines:
        rect(s, Inches(0.9), Emu(y + Inches(0.05)), Inches(0.13), Inches(0.13), color)
        text(s, Inches(1.25), y, Inches(10.9), Inches(1.0),
             [[(body, SANS, 15, INK, False, False)]])
        y = Emu(y + Inches(1.05))
    return s


def frame_slide(prs, f, index, total):
    s = blank(prs)
    rect(s, 0, 0, W, H, PAPER)
    # Left photo mat.
    mat_w = Inches(7.5)
    rect(s, 0, 0, mat_w, H, FOREST if f["kind"] == "ready" else RGBColor(0x3a, 0x3a, 0x3a))
    buf, wp, hp = load_img(f)
    iw, ih = fit(wp, hp, Inches(6.7), Inches(6.5))
    ix = Emu(int((mat_w - iw) / 2))
    iy = Emu(int((H - ih) / 2))
    s.shapes.add_picture(buf, ix, iy, iw, ih)
    # Year tab.
    tab = rect(s, 0, Inches(0.0), Inches(2.0), Inches(0.62), GOLD if f["kind"]=="ready" else RGBColor(0x8a,0x8a,0x8a))
    tf = tab.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(10)
    p = tf.paragraphs[0]; r = p.add_run(); r.text = str(f["year_display"])
    r.font.name = SERIF; r.font.size = Pt(18); r.font.bold = True
    r.font.color.rgb = INK if f["kind"]=="ready" else WHITE
    # Right text column.
    cx, cw = Inches(7.9), Inches(4.9)
    text(s, cx, Inches(0.7), cw, Inches(0.4),
         [[(f"FRAME {index} OF {total}", SANS, 12, GOLD, True, False)]])
    text(s, cx, Inches(1.15), cw, Inches(1.2),
         [[(f["title"], SERIF, 28, INK, False, False)]])
    # Visitor-facing caption — museum-label voice, the main body text on the wall.
    text(s, cx, Inches(2.5), cw, Inches(3.4),
         [[(f.get("visitor_caption") or f.get("caption", ""), SERIF, 16, RGBColor(0x2c,0x36,0x30), False, False)]])
    text(s, cx, Inches(6.0), cw, Inches(0.4),
         [[("Credit: " + (f["credit"] or "Unknown"), SANS, 11, MUTED, False, False)]])
    # Status footer.
    if f["kind"] == "ready":
        chip(s, cx, Inches(6.55), f'READY · {f["long_edge_in_150"]}"  @150 PPI', READY)
    else:
        w = chip(s, cx, Inches(6.35), "TO SOURCE", RGBColor(0x5b,0x5b,0x5b))
        chip(s, Emu(cx + w + Inches(0.15)), Inches(6.35),
             f'TRACK {f["track"]}', TRACK_COLOR.get(f["track"], MUTED))
        text(s, cx, Inches(6.85), cw, Inches(0.5),
             [[(f["source"] or "", SANS, 11, MUTED, False, False)]])
    # Speaker notes: the internal case for the image (why it earns a place), plus
    # sourcing status for gaps. Not shown on the slide — the presenter's talking points.
    argument = f.get("argument") or f.get("rationale") or ""
    note = f"WHY THIS IMAGE IS ON THE WALL\n{argument}\n"
    if f["kind"] == "ready":
        note += (f"\nProduction status: ready to print — {f['long_edge_in_150']}\" long "
                 f"edge at 150 PPI ({f['pixels']} px). Taste-matched stand-in for "
                 f"Kelly's curated vision.")
    else:
        note += (f"\nSourcing status: TO RE-ACQUIRE — no printable version exists. "
                 f"{f.get('track_label') or ('Track ' + str(f['track']))}. "
                 f"Source lead: {f.get('source') or '—'}.")
    note += "\n\n(Caption is draft visitor-facing text; confirm facts and rights before public use.)"
    s.notes_slide.notes_text_frame.text = note
    return s


def ask_slide(prs, data):
    s = blank(prs)
    rect(s, 0, 0, W, H, FOREST)
    rect(s, 0, 0, W, Inches(0.14), GOLD)
    text(s, Inches(0.9), Inches(0.75), Inches(11.5), Inches(0.9),
         [[("What we need to finish the wall", SERIF, 36, WHITE, False, False)]])
    gaps = [f for f in data["frames"] if f["kind"] == "gap"]
    by = {"A": [], "B": [], "C": []}
    for g in gaps:
        by.get(g["track"], by.setdefault(g["track"], [])).append(g)
    blocks = [
        ("A", "Track A · Online", "Calisphere coll. 1828/1582, OAC, Record Gazette. "
         "Start with the 1875 Depot — it's already online with an ARK.", by.get("A", [])),
        ("B", "Track B · Malki + Tribe", "Malki Museum sourcing AND Morongo Tribe "
         "cultural consent before any public use of the Native-heritage images.", by.get("B", [])),
        ("C", "Track C · Storage", "The unindexed Library storage. Bring the printed "
         "pull-list; triage on the first visit.", by.get("C", [])),
    ]
    y = Inches(1.9)
    for track, head, body, items in blocks:
        rect(s, Inches(0.9), y, Inches(0.13), Inches(1.4), TRACK_COLOR[track])
        titles = "  ·  ".join(g["title"].split(" From ")[0][:34] for g in items)
        text(s, Inches(1.25), y, Inches(11.0), Inches(1.5),
             [[(f"{head}   ({len(items)})", SANS, 17, GOLD, True, False)],
              [(body, SANS, 13, RGBColor(0xCF, 0xDA, 0xD3), False, False)],
              [(titles, SERIF, 12.5, WHITE, False, True)]])
        y = Emu(y + Inches(1.6))
    return s


def _span(data):
    years = [f["year"] for f in data["frames"] if f["year"] is not None]
    return f"{min(years)}–{max(years)}" if years else "—"


def main():
    data = json.loads(DATA.read_text(encoding="utf-8"))
    frames = data["frames"]
    prs = Presentation()
    prs.slide_width, prs.slide_height = W, H
    title_slide(prs, data)
    howto_slide(prs)
    for i, f in enumerate(frames, 1):
        frame_slide(prs, f, i, len(frames))
    ask_slide(prs, data)
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    kb = OUT.stat().st_size / 1024
    print(f"Wrote {OUT.relative_to(ROOT)} ({kb:.0f} KB, {len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
